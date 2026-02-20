"""
Generate a PowerPoint deck for the MemoryCharm Claim Flow UX review.
Each slide includes a breadcrumb "You Are Here" progress diagram,
an image placeholder, and a plain-language description.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────────────────────────
BG_COLOR       = RGBColor(0x1A, 0x1A, 0x2E)   # deep indigo background
ACCENT_GOLD    = RGBColor(0xD4, 0xAF, 0x37)   # warm gold
TEXT_LIGHT     = RGBColor(0xF5, 0xF0, 0xE1)   # parchment white
TEXT_DIM       = RGBColor(0x88, 0x88, 0x99)   # muted label
STEP_INACTIVE  = RGBColor(0x33, 0x33, 0x50)   # dim circle
STEP_ACTIVE_BG = RGBColor(0xD4, 0xAF, 0x37)   # gold highlight
STEP_DONE_BG   = RGBColor(0x3A, 0x6B, 0x4C)   # muted green for completed
CONNECTOR_DIM  = RGBColor(0x44, 0x44, 0x60)
CONNECTOR_DONE = RGBColor(0x3A, 0x6B, 0x4C)
PLACEHOLDER_BG = RGBColor(0x22, 0x22, 0x3A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
BLACK          = RGBColor(0x00, 0x00, 0x00)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Step definitions ─────────────────────────────────────────────────────
# We show the FULL possible flow (new user + glyph lock = 6 steps).
# For the "open charm" variant we'll add a note on the protection slide.

STEPS = [
    {
        "id": "register",
        "short": "Register",
        "title": "Step 1 \u2014 Register as a Keeper",
        "description": (
            "Before you can claim a charm, we need to know who you are.\n\n"
            "On this screen you'll fill in your name, address, email, and "
            "phone number. Your email is pre-filled from your sign-in.\n\n"
            "You'll also need to accept the Terms of Passage before continuing.\n\n"
            "If you've already registered in a previous session, this step "
            "is skipped automatically."
        ),
    },
    {
        "id": "memory_type",
        "short": "Memory Type",
        "title": "Step 2 \u2014 Choose Your Memory Type",
        "description": (
            "Every charm holds one memory. Here you pick what kind:\n\n"
            "\u2022  Video \u2013 a moving picture, captured in time\n"
            "\u2022  Image \u2013 a still photograph (you can add more than one)\n"
            "\u2022  Audio \u2013 a voice recording or melody\n\n"
            "Tap the one that matches the file you'd like to upload, "
            "then press Continue."
        ),
    },
    {
        "id": "upload",
        "short": "Upload",
        "title": "Step 3 \u2014 Upload Your Memory",
        "description": (
            "Now you'll select the actual file from your device.\n\n"
            "The screen shows a file picker \u2014 tap it, find your video, "
            "photo(s), or audio file, and select it. You'll see the file "
            "name and size appear on screen.\n\n"
            "Files can be up to 150 MB. Once you've chosen your file, "
            "press Continue."
        ),
    },
    {
        "id": "protection",
        "short": "Protection",
        "title": "Step 4 \u2014 Choose How to Protect Your Charm",
        "description": (
            "You decide who can see this memory:\n\n"
            "\u2022  Open \u2013 anyone who scans the charm can view it immediately.\n"
            "\u2022  Glyph Lock \u2013 the viewer must pick the correct secret "
            "symbol before the memory is revealed.\n\n"
            "If you choose Open, the charm is sealed right away and you'll "
            "jump to the final confirmation.\n\n"
            "If you choose Glyph Lock, you'll pick your secret symbol "
            "on the next screen."
        ),
    },
    {
        "id": "glyph",
        "short": "Secret Glyph",
        "title": "Step 5 \u2014 Pick Your Secret Glyph",
        "description": (
            "This screen shows a grid of 18 symbols \u2014 stars, hearts, "
            "moons, butterflies, and more.\n\n"
            "Tap the one you want as your charm's secret key. Anyone who "
            "later scans the charm will be shown 9 random symbols and "
            "must pick the right one to unlock the memory.\n\n"
            "They get 3 attempts. Choose something meaningful to you "
            "that the right person will recognise.\n\n"
            "This step only appears if you chose Glyph Lock."
        ),
    },
    {
        "id": "done",
        "short": "Sealed!",
        "title": "Step 6 \u2014 Your Charm Is Sealed",
        "description": (
            "You're done! The memory has been securely bound to your charm.\n\n"
            "This screen shows a preview of what you uploaded \u2014 you can "
            "play a video or audio clip, or scroll through your photos.\n\n"
            "Press \"View Charm\" to see exactly what someone will "
            "experience when they scan the physical charm."
        ),
    },
]

# ── Helpers ───────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set a solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Make corners fairly round
    shape.adjustments[0] = 0.25
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=TEXT_LIGHT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    """Add a text box with a single run of styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def draw_breadcrumb(slide, steps, active_index, top_y):
    """
    Draw a horizontal breadcrumb bar across the slide.
    Completed steps are green, the active step is gold, future steps are dim.
    Connectors link the circles.
    """
    n = len(steps)
    circle_d = Inches(0.55)
    gap = Inches(1.5)          # center-to-center spacing
    total_w = (n - 1) * gap
    start_x = int((SLIDE_W - total_w) / 2 - circle_d / 2)

    connector_h = Pt(3)
    connector_top = int(top_y + circle_d / 2 - connector_h / 2)

    for i in range(n):
        cx = int(start_x + i * gap)

        # ── connector line to previous circle ──
        if i > 0:
            line_left = int(start_x + (i - 1) * gap + circle_d)
            line_w = int(gap - circle_d)
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, line_left, connector_top, line_w, connector_h
            )
            conn.fill.solid()
            conn.line.fill.background()
            if i <= active_index:
                conn.fill.fore_color.rgb = CONNECTOR_DONE
            else:
                conn.fill.fore_color.rgb = CONNECTOR_DIM

        # ── circle ──
        if i < active_index:
            bg = STEP_DONE_BG
            fg = WHITE
        elif i == active_index:
            bg = STEP_ACTIVE_BG
            fg = BLACK
        else:
            bg = STEP_INACTIVE
            fg = TEXT_DIM

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx, top_y, circle_d, circle_d
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg
        circle.line.fill.background()

        # Number inside circle
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        if i < active_index:
            run.text = "\u2713"   # checkmark
        else:
            run.text = str(i + 1)
        run.font.size = Pt(16)
        run.font.color.rgb = fg
        run.font.bold = True
        run.font.name = "Segoe UI"

        # Label below
        label_top = int(top_y + circle_d + Inches(0.08))
        lbl = add_textbox(
            slide, cx - Inches(0.4), label_top, Inches(1.35), Inches(0.35),
            steps[i]["short"],
            font_size=9,
            color=ACCENT_GOLD if i == active_index else TEXT_DIM,
            bold=(i == active_index),
            alignment=PP_ALIGN.CENTER,
        )

    # "You are here" callout above active circle
    arrow_x = int(start_x + active_index * gap - Inches(0.25))
    arrow_y = int(top_y - Inches(0.55))
    add_textbox(
        slide, arrow_x, arrow_y, Inches(1.1), Inches(0.35),
        "\u25BC  You are here",
        font_size=10, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER,
    )


def draw_breadcrumb_compact(slide, steps, active_index, top_y, region_left, region_w):
    """
    Draw a breadcrumb bar constrained within a specific horizontal region.
    Used for the right-column layout on step slides.
    """
    n = len(steps)
    circle_d = Inches(0.45)
    gap = Inches(1.25)          # center-to-center spacing
    total_w = (n - 1) * gap
    start_x = int(region_left + (region_w - total_w) / 2 - circle_d / 2)

    connector_h = Pt(3)
    connector_top = int(top_y + circle_d / 2 - connector_h / 2)

    for i in range(n):
        cx = int(start_x + i * gap)

        # ── connector line to previous circle ──
        if i > 0:
            line_left = int(start_x + (i - 1) * gap + circle_d)
            line_w = int(gap - circle_d)
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, line_left, connector_top, line_w, connector_h
            )
            conn.fill.solid()
            conn.line.fill.background()
            if i <= active_index:
                conn.fill.fore_color.rgb = CONNECTOR_DONE
            else:
                conn.fill.fore_color.rgb = CONNECTOR_DIM

        # ── circle ──
        if i < active_index:
            bg = STEP_DONE_BG
            fg = WHITE
        elif i == active_index:
            bg = STEP_ACTIVE_BG
            fg = BLACK
        else:
            bg = STEP_INACTIVE
            fg = TEXT_DIM

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx, top_y, circle_d, circle_d
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg
        circle.line.fill.background()

        # Number inside circle
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        if i < active_index:
            run.text = "\u2713"
        else:
            run.text = str(i + 1)
        run.font.size = Pt(13)
        run.font.color.rgb = fg
        run.font.bold = True
        run.font.name = "Segoe UI"

        # Label below
        label_top = int(top_y + circle_d + Inches(0.05))
        add_textbox(
            slide, cx - Inches(0.35), label_top, Inches(1.15), Inches(0.3),
            steps[i]["short"],
            font_size=8,
            color=ACCENT_GOLD if i == active_index else TEXT_DIM,
            bold=(i == active_index),
            alignment=PP_ALIGN.CENTER,
        )

    # "You are here" callout above active circle
    arrow_x = int(start_x + active_index * gap - Inches(0.2))
    arrow_y = int(top_y - Inches(0.45))
    add_textbox(
        slide, arrow_x, arrow_y, Inches(0.9), Inches(0.3),
        "\u25BC  You are here",
        font_size=9, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER,
    )


# ── Build presentation ───────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  TITLE SLIDE                                                        ║
# ╚══════════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)

add_textbox(
    slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1),
    "MemoryCharm", font_size=48, color=ACCENT_GOLD, bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.6),
    "Claim Flow \u2014 UX Walkthrough", font_size=28, color=TEXT_LIGHT,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(1), Inches(4.2), Inches(11.3), Inches(1),
    (
        "This deck walks through every screen a user sees when they scan a charm\n"
        "and bind a memory to it. Each slide shows where the user is in the\n"
        "overall journey and describes what they do on that screen."
    ),
    font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER,
)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  OVERVIEW SLIDE                                                      ║
# ╚══════════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)

add_textbox(
    slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.6),
    "The Journey at a Glance", font_size=32, color=ACCENT_GOLD, bold=True,
    alignment=PP_ALIGN.CENTER,
)

# Draw the breadcrumb with no active step (all neutral)
# We'll draw it manually as a "map"
overview_top = Inches(1.5)
n = len(STEPS)
circle_d = Inches(0.7)
gap = Inches(1.7)
total_w = (n - 1) * gap
start_x = int((SLIDE_W - total_w) / 2 - circle_d / 2)

for i in range(n):
    cx = int(start_x + i * gap)
    # connector
    if i > 0:
        line_left = int(start_x + (i - 1) * gap + circle_d)
        line_w = int(gap - circle_d)
        conn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, line_left,
            int(overview_top + circle_d / 2 - Pt(2)),
            line_w, Pt(4)
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = ACCENT_GOLD
        conn.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx, overview_top, circle_d, circle_d
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i + 1)
    run.font.size = Pt(22)
    run.font.color.rgb = BLACK
    run.font.bold = True
    run.font.name = "Segoe UI"

    # label + mini description
    label_top = int(overview_top + circle_d + Inches(0.15))
    add_textbox(
        slide, cx - Inches(0.5), label_top, Inches(1.7), Inches(0.3),
        STEPS[i]["short"], font_size=13, color=TEXT_LIGHT, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

# Add two flow variant notes at the bottom
add_textbox(
    slide, Inches(1), Inches(3.6), Inches(11.3), Inches(1.2),
    (
        "Two possible paths through the flow:\n\n"
        "\u2022  Open charm (no lock):  Register \u2192 Memory Type \u2192 Upload "
        "\u2192 Protection \u2192 Sealed!   (steps 1\u20134, then 6)\n"
        "\u2022  Glyph-locked charm:    Register \u2192 Memory Type \u2192 Upload "
        "\u2192 Protection \u2192 Secret Glyph \u2192 Sealed!   (all 6 steps)\n\n"
        "If the user has already registered, step 1 is skipped."
    ),
    font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.LEFT,
)

# Decorative divider
div = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(5.2), Inches(9.3), Pt(1)
)
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_GOLD
div.line.fill.background()

add_textbox(
    slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.8),
    (
        "The following slides show each step in detail. "
        "Add a screenshot to the placeholder on each slide for the full picture."
    ),
    font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER,
)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  STEP SLIDES                                                        ║
# ╚══════════════════════════════════════════════════════════════════════╝

for idx, step in enumerate(STEPS):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_COLOR)

    # ── Layout: LEFT 1/3 = screenshot, RIGHT 2/3 = breadcrumb + text ──
    margin = Inches(0.5)
    col_gap = Inches(0.4)

    # Left column: screenshot placeholder (full height, portrait proportion)
    ph_left = margin
    ph_top = margin
    ph_w = Inches(3.9)
    ph_h = SLIDE_H - 2 * margin  # nearly full slide height

    ph_shape = add_rounded_rect(
        slide, ph_left, ph_top, ph_w, ph_h,
        PLACEHOLDER_BG, line_color=ACCENT_GOLD,
    )

    # "Insert screenshot here" label centred in placeholder
    add_textbox(
        slide, ph_left, int(ph_top + ph_h / 2 - Inches(0.5)),
        ph_w, Inches(0.6),
        "\u2702  Insert screenshot here",
        font_size=16, color=TEXT_DIM, alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, ph_left, int(ph_top + ph_h / 2 + Inches(0.1)),
        ph_w, Inches(0.5),
        "(right-click \u2192 Change Picture)",
        font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER,
    )

    # Right column starts after the screenshot + gap
    right_left = int(ph_left + ph_w + col_gap)
    right_w = int(SLIDE_W - right_left - margin)

    # ── Step title (top of right column) ──
    title_top = Inches(0.5)
    add_textbox(
        slide, right_left, title_top, right_w, Inches(0.55),
        step["title"],
        font_size=26, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.LEFT,
    )

    # ── Breadcrumb (below title) ──
    breadcrumb_top = Inches(1.4)
    draw_breadcrumb_compact(slide, STEPS, idx, breadcrumb_top, right_left, right_w)

    # ── Description panel (below breadcrumb, fills remaining height) ──
    desc_top = Inches(2.8)
    desc_h = int(SLIDE_H - desc_top - margin)

    panel = add_rounded_rect(
        slide, right_left, desc_top, right_w, desc_h,
        RGBColor(0x20, 0x20, 0x38), line_color=None,
    )

    # "What the user does here:" header
    add_textbox(
        slide, int(right_left + Inches(0.3)), int(desc_top + Inches(0.25)),
        int(right_w - Inches(0.6)), Inches(0.4),
        "What the user does here:",
        font_size=14, color=ACCENT_GOLD, bold=True,
    )

    # Description body
    add_textbox(
        slide, int(right_left + Inches(0.3)), int(desc_top + Inches(0.7)),
        int(right_w - Inches(0.6)), int(desc_h - Inches(1.0)),
        step["description"],
        font_size=13, color=TEXT_LIGHT,
    )

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  END SLIDE                                                          ║
# ╚══════════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)

add_textbox(
    slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1),
    "Thank you!", font_size=44, color=ACCENT_GOLD, bold=True,
    alignment=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2),
    (
        "We'd love your feedback on this flow.\n\n"
        "Is each step clear?  Anything confusing or unnecessary?\n"
        "Would you change the order of any steps?"
    ),
    font_size=18, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER,
)

# ── Save ──────────────────────────────────────────────────────────────
output_path = r"c:\Users\appli\source\repos\MemoryCharm\MemoryCharm_ClaimFlow_UX_v2.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
